"""Answer -> a DESIGNED PowerPoint deck, in the spirit of Coreworks' output (themed layouts, a
big headline number per slide) with the one thing a generic generator cannot promise: every figure
carries the source cell it came from.

Still a renderer, not a rewrite (D40): it never computes or formats a number. Every figure is
`Value.formatted`, already made by execute(). The design lives here; the numbers do not.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

import re
import sys
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from answer import Answer, Value                       # noqa: E402
from present import pretty_sql                         # noqa: E402

# ── theme ────────────────────────────────────────────────────────────────────────
ORANGE = RGBColor(0xE8, 0x5D, 0x2A)
INK    = RGBColor(0x19, 0x19, 0x19)
MUTED  = RGBColor(0x8A, 0x8A, 0x8A)
CARD   = RGBColor(0xF5, 0xF2, 0xEE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x1B, 0x7F, 0x3B)
AMBER  = RGBColor(0xB4, 0x69, 0x0E)
BLUE   = RGBColor(0x1F, 0x5C, 0xA8)
FONT   = "Arial"

STATUS = {"answered": ("ANSWERED", GREEN), "clarify": ("NEEDS INPUT", AMBER),
          "abstained": ("ABSTAINED", BLUE)}
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _clean(x) -> str:
    return " ".join(str(x).split())


def _text(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _run(para, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT, mono=False):
    para.alignment = align
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Consolas" if mono else FONT
    return r


def _rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _chip(slide, l, t, label, color):
    sp = _rect(slide, l, t, 1.9, 0.42, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame; tf.word_wrap = False
    _run(tf.paragraphs[0], label, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _footnote(slide, cites):
    """Citations, demoted to a thin line at the foot of the slide. The number and the story are the
    slide; the cells are the receipt underneath, there if you want to check but not shouting."""
    uniq, formula = [], None
    for c in cites:
        tag = f"{c.sheet}!{c.a1}"
        if tag not in uniq:
            uniq.append(tag)
        if c.formula and not formula:
            formula = c.formula
    if not uniq:
        return
    shown = ", ".join(uniq[:3]) + (f"  +{len(uniq) - 3} more" if len(uniq) > 3 else "")
    line = "Source: " + shown + (f"   ·  {formula}" if formula else "")
    tf = _text(slide, 0.9, 7.02, 11.6, 0.4)
    _run(tf.paragraphs[0], line[:140], 9, color=MUTED, mono=True)


def _headline(answer: Answer):
    """The numeric star of the slide, and any label that goes with it (the winning entity)."""
    numeric, label = None, None
    for v in answer.slots.values():
        is_num = isinstance(v.raw, (int, float)) or re.fullmatch(r"-?[\d,]+\.?\d*", str(v.formatted))
        if is_num and numeric is None:
            numeric = v
        elif not is_num and label is None:
            label = v
    return numeric, label


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def cover(prs, title, subtitle, n):
    s = _blank(prs)
    _rect(s, 0, 0, 0.35, 7.5, ORANGE)
    tf = _text(s, 0.9, 2.3, 11.5, 2.5)
    _run(tf.paragraphs[0], title, 44, bold=True, color=INK)
    _run(tf.add_paragraph(), subtitle, 18, color=MUTED)
    p = tf.add_paragraph(); p.space_before = Pt(18)
    _run(p, f"{n} findings · every number traces to a source cell", 13, color=ORANGE, bold=True)


def stat_slide(prs, answer: Answer):
    s = _blank(prs)
    word, color = STATUS[answer.status]
    _rect(s, 0, 0, 13.333, 0.14, color)              # top accent rule
    _chip(s, 0.9, 0.5, word, color)

    # question as eyebrow
    tf = _text(s, 0.9, 1.15, 11.5, 1.0)
    _run(tf.paragraphs[0], answer.question, 20, bold=True, color=INK)

    if answer.status == "answered":
        numeric, label = _headline(answer)
        big = numeric.formatted if numeric else answer.text()
        unit = (numeric.unit if numeric and numeric.unit else "")
        # giant number
        nt = _text(s, 0.9, 2.5, 11.6, 2.0, anchor=MSO_ANCHOR.MIDDLE)
        _run(nt.paragraphs[0], big, 72, bold=True, color=INK)
        if unit or label:
            sub = nt.add_paragraph()
            _run(sub, "  ".join(x for x in [_clean(label.formatted) if label else "", unit] if x),
                 18, color=MUTED)
        # the sentence — now full width, the story is the slide
        st = _text(s, 0.9, 4.9, 11.6, 1.6)
        _run(st.paragraphs[0], answer.text(), 20, color=INK)
        _footnote(s, answer.all_citations())
    else:
        # refusal / clarify: the message is the content, and there is NO number
        bt = _text(s, 0.9, 2.6, 11.0, 2.0)
        _run(bt.paragraphs[0], answer.text(), 22, color=INK)
        if answer.scope_options:
            op = _text(s, 0.9, 4.6, 11.0, 2.0)
            _run(op.paragraphs[0], "Options", 13, bold=True, color=MUTED)
            for o in answer.scope_options:
                _run(op.add_paragraph(), f"›  {_clean(o)}", 15, color=INK)


def _evidence(v: Value, depth: int = 0) -> list[str]:
    pad = "  " * depth
    out = []
    for c in v.citations:
        f = f"  [{c.formula}]" if c.formula else ""
        out.append(f"{pad}{c.sheet}!{c.a1} = {c.raw_value}{f}")
    for name, part in v.parts.items():
        out.append(f"{pad}{name} = {part.formatted}")
        out += _evidence(part, depth + 1)
    return out


def closing_slide(prs, text):
    s = _blank(prs)
    _rect(s, 0, 0, 0.35, 7.5, ORANGE)
    tf = _text(s, 0.9, 3.0, 11.5, 2.0, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], "In short", 14, bold=True, color=ORANGE)
    _run(tf.add_paragraph(), text, 26, bold=True, color=INK)


def leaderboard_slide(prs, lb):
    """A top-N finding: a horizontal bar chart, with the source cells listed beside it."""
    s = _blank(prs)
    _rect(s, 0, 0, 13.333, 0.14, GREEN)
    _chip(s, 0.9, 0.5, "ANSWERED", GREEN)
    tf = _text(s, 0.9, 1.15, 11.5, 0.9)
    _run(tf.paragraphs[0], lb.question, 20, bold=True, color=INK)

    # bar chart (highest at top -> feed reversed, since the chart plots first category at bottom)
    rows = list(reversed(lb.rows))
    data = CategoryChartData()
    data.categories = [_clean(e)[:28] for e, _f, _r, _c in rows]
    data.add_series(_clean(lb.measure)[:30], [float(r) for _e, _f, r, _c in rows])
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(2.2),
                            Inches(11.5), Inches(4.6), data)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    plot = ch.plots[0]
    plot.has_data_labels = True
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = ORANGE
    _footnote(s, lb.all_citations())


def trend_slide(prs, tr):
    """A time series as a line chart, with the direction-of-travel headline above it."""
    s = _blank(prs)
    _rect(s, 0, 0, 13.333, 0.14, GREEN)
    _chip(s, 0.9, 0.5, "TREND", GREEN)
    tf = _text(s, 0.9, 1.15, 11.6, 1.4)
    _run(tf.paragraphs[0], tr.headline or tr.text(), 22, bold=True, color=INK)

    data = CategoryChartData()
    data.categories = [_clean(pl)[:12] for pl, _f, _r, _c in tr.rows]
    data.add_series(_clean(tr.measure)[:30], [float(r) for _p, _f, r, _c in tr.rows])
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.9), Inches(2.7),
                            Inches(11.5), Inches(4.1), data)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    ch.series[0].format.line.color.rgb = ORANGE
    _footnote(s, tr.all_citations())


def build_deck(answers: list[Answer], path: Path, title="Data analysis", subtitle="",
               closing="") -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    cover(prs, title, subtitle, len(answers))
    for a in answers:
        if hasattr(a, "headline") and hasattr(a, "rows"):    # a Trend finding
            trend_slide(prs, a)
        elif hasattr(a, "rows") and hasattr(a, "measure"):   # a Leaderboard finding
            leaderboard_slide(prs, a)
        else:
            stat_slide(prs, a)
    if closing:
        closing_slide(prs, closing)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
